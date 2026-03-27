from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

app = Flask(__name__)

@app.route("/")
def home():
    return render_template("index.html")


@app.route("/submit", methods=["POST"])
def submit():

    data = request.form

    prs = Presentation()

    # -------- BULLET FORMAT FUNCTION --------
    def format_bullets(text):
        if not text:
            return ""
        sentences = text.split(". ")
        formatted = ""
        for sentence in sentences:
            sentence = sentence.strip()
            if sentence:
                formatted += "• " + sentence + "\n"
        return formatted

    # -------- TITLE SLIDE (PREMIUM STYLE) --------
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(2, 6, 23)

    startup = data.get("startup", "")
    tagline = data.get("tagline", "")

    slide.shapes.title.text = startup

    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 191, 255)

    subtitle = slide.placeholders[1]
    subtitle.text = tagline + "\n\nPowered by InvestoDeck"

    sub_para = subtitle.text_frame.paragraphs[0]
    sub_para.font.size = Pt(20)
    sub_para.font.color.rgb = RGBColor(255, 255, 255)

    # -------- FUNCTION FOR CONTENT SLIDES --------
    def add_slide(title, content):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)

        # Title
        slide.shapes.title.text = title
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.size = Pt(34)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 191, 255)

        # Content
        content_shape = slide.placeholders[1]
        content_shape.text = content

        for para in content_shape.text_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(255, 255, 255)

        # Footer
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(5)
        height = Inches(0.3)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = "InvestoDeck | AI Pitch Generator"

        textbox.text_frame.paragraphs[0].font.size = Pt(12)
        textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)

    # -------- CONTENT SLIDES --------
    add_slide("Problem Statement", format_bullets(data.get("problem", "")))
    add_slide("Solution", format_bullets(data.get("solution", "")))
    add_slide("Target Market", format_bullets(data.get("market", "")))
    add_slide("Business Model", format_bullets(data.get("revenue", "")))
    add_slide("Funding Requirement", format_bullets(data.get("funding", "")))

    # -------- FINANCIAL CHART (IMPROVED) --------
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = "Financial Projections"
    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.size = Pt(34)
    title_para.font.color.rgb = RGBColor(0, 191, 255)

    chart_data = ChartData()
    chart_data.categories = ['Year 1', 'Year 2', 'Year 3']

    year1 = int(data.get("year1", 0) or 0)
    year2 = int(data.get("year2", 0) or 0)
    year3 = int(data.get("year3", 0) or 0)

    chart_data.add_series('Revenue', (year1, year2, year3))

    x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(4)

    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        x, y, cx, cy,
        chart_data
    )

    filename = "Investor_Pitch_Deck.pptx"
    prs.save(filename)

    return send_file(filename, as_attachment=True)


if __name__ == "__main__":
    app.run(debug=True)
